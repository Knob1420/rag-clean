"""
文档格式转换器 — 统一入口：任意格式 → Markdown 文本

支持的格式及降级策略：
├── .md    → 直接读取
├── .pdf   → MinerU（复用 parser.py）
├── .doc   → antiword → libreoffice 降级
├── .docx  → python-docx → XML 解析 → libreoffice 三级降级
├── .pptx  → python-pptx
└── .ppt   → 暂不支持
"""

import hashlib
import subprocess
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional

from loguru import logger

# 支持的文件格式
SUPPORTED_FORMATS = {".pdf", ".doc", ".docx", ".pptx", ".ppt", ".md"}


def detect_format(file_path: Path) -> str:
    """返回格式标识: pdf/doc/docx/pptx/ppt/md/unknown"""
    ext = Path(file_path).suffix.lower()
    format_map = {
        ".pdf": "pdf",
        ".doc": "doc",
        ".docx": "docx",
        ".pptx": "pptx",
        ".ppt": "ppt",
        ".md": "md",
    }
    return format_map.get(ext, "unknown")


def convert_to_markdown(file_path: str) -> str:
    """
    统一入口：任意格式 → Markdown 文本

    Args:
        file_path: 文件路径

    Returns:
        Markdown 文本内容

    Raises:
        ValueError: 不支持的格式
        RuntimeError: 转换失败
    """
    path = Path(file_path)
    fmt = detect_format(path)

    if fmt == "unknown":
        raise ValueError(f"不支持的文件格式: {path.suffix}")

    if fmt == "md":
        return path.read_text(encoding="utf-8")

    if fmt == "pdf":
        return _convert_pdf(path)

    if fmt == "doc":
        return _convert_doc(path)

    if fmt == "docx":
        return _convert_docx(path)

    if fmt == "pptx":
        return _convert_pptx(path)

    if fmt == "ppt":
        raise RuntimeError("老 PPT 格式暂不支持，请转换为 PPTX")

    raise ValueError(f"未知的格式标识: {fmt}")


# ── 内部转换函数 ───────────────────────────────────────


def _convert_pdf(path: Path) -> str:
    """PDF → Markdown（复用 parser.py 的 MinerUPDFProcessor）"""
    from core.ingestion.parser import MinerUPDFProcessor

    processor = MinerUPDFProcessor()

    # 检查缓存目录
    cache_dir = _converters_cache_dir() / f"{path.stem}"
    md_file = cache_dir / f"{path.stem}.md"
    if md_file.exists():
        logger.info(f"[Converter] PDF 缓存命中: {md_file}")
        return md_file.read_text(encoding="utf-8")

    md_content = processor.process_pdf(str(path))
    cache_dir.mkdir(parents=True, exist_ok=True)
    md_file.write_text(md_content, encoding="utf-8")
    return md_content


def _convert_doc(path: Path) -> str:
    """DOC → Markdown: antiword → libreoffice 降级"""
    # 检查缓存
    cached = _load_cache(path)
    if cached is not None:
        return cached

    # 方法 1: antiword
    result = _try_antiword(path)
    if result is not None:
        _save_cache(path, result)
        return result

    # 方法 2: libreoffice
    result = _convert_with_libreoffice(path)
    if result is not None:
        _save_cache(path, result)
        return result

    raise RuntimeError(f"DOC 转换失败（antiword 和 libreoffice 均不可用）: {path.name}")


def _convert_docx(path: Path) -> str:
    """DOCX → Markdown: python-docx → XML 解析 → libreoffice 三级降级"""
    # 检查缓存
    cached = _load_cache(path)
    if cached is not None:
        return cached

    # 方法 1: python-docx
    result = _try_python_docx(path)
    if result is not None:
        _save_cache(path, result)
        return result

    # 方法 2: 直接解析 XML
    result = _parse_docx_xml(path)
    if result is not None:
        _save_cache(path, result)
        return result

    # 方法 3: libreoffice
    result = _convert_with_libreoffice(path)
    if result is not None:
        _save_cache(path, result)
        return result

    raise RuntimeError(f"DOCX 转换失败（所有方法均不可用）: {path.name}")


def _convert_pptx(path: Path) -> str:
    """PPTX → Markdown: python-pptx"""
    # 检查缓存
    cached = _load_cache(path)
    if cached is not None:
        return cached

    result = _try_python_pptx(path)
    if result is not None:
        _save_cache(path, result)
        return result

    raise RuntimeError(f"PPTX 转换失败: {path.name}")


# ── 具体转换实现 ───────────────────────────────────────


def _try_antiword(path: Path) -> Optional[str]:
    """使用 antiword 转换 DOC 文件"""
    try:
        result = subprocess.run(
            ["antiword", "-m", "UTF-8", str(path)],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            md_lines = [f"# {path.stem}\n"]
            md_lines.append(result.stdout)
            return "\n".join(md_lines)
        logger.warning(f"antiword 转换失败: {result.stderr}")
    except FileNotFoundError:
        logger.info("antiword 未安装，跳过")
    except subprocess.TimeoutExpired:
        logger.warning("antiword 转换超时")
    except Exception as e:
        logger.warning(f"antiword 异常: {e}")
    return None


def _try_python_docx(path: Path) -> Optional[str]:
    """使用 python-docx 转换 DOCX 文件"""
    try:
        from docx import Document

        doc = Document(path)
        md_lines = [f"# {path.stem}\n"]

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                md_lines.append("")
                continue

            if para.style.name.startswith("Heading"):
                try:
                    level = int(para.style.name.replace("Heading ", "").strip())
                    md_lines.append(f"{'#' * (level + 1)} {text}")
                except ValueError:
                    md_lines.append(f"## {text}")
            else:
                md_lines.append(text)

        return "\n".join(md_lines)

    except ImportError:
        logger.info("python-docx 未安装，跳过")
        return None
    except Exception as e:
        logger.warning(f"python-docx 转换失败: {e}")
        return None


def _parse_docx_xml(path: Path) -> Optional[str]:
    """直接解析 word/document.xml（处理非标准 DOCX）"""
    try:
        with zipfile.ZipFile(path, "r") as zf:
            doc_xml = None
            for doc_path in ["word/document.xml", "word\\document.xml"]:
                try:
                    doc_xml = zf.read(doc_path)
                    break
                except KeyError:
                    continue

            if doc_xml is None:
                return None

            root = ET.fromstring(doc_xml)
            namespaces = {
                "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            }

            md_lines = [f"# {path.stem}\n"]
            for para in root.findall(".//w:p", namespaces):
                texts = [t.text or "" for t in para.findall(".//w:t", namespaces)]
                para_text = "".join(texts)
                if para_text.strip():
                    md_lines.append(para_text)
                else:
                    md_lines.append("")

            return "\n".join(md_lines)

    except Exception as e:
        logger.warning(f"XML 解析失败: {e}")
        return None


def _convert_with_libreoffice(path: Path) -> Optional[str]:
    """使用 libreoffice 将 DOC/DOCX 转换为纯文本"""
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)

            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "txt:Text (encoded)",
                    "--outdir",
                    str(temp_path),
                    str(path),
                ],
                capture_output=True,
                text=True,
                timeout=60,
            )

            if result.returncode != 0:
                logger.error(f"libreoffice 转换失败: {result.stderr}")
                return None

            txt_files = list(temp_path.glob("*.txt"))
            if not txt_files:
                logger.error("libreoffice 未生成输出文件")
                return None

            text_content = txt_files[0].read_text(encoding="utf-8", errors="ignore")
            md_lines = [f"# {path.stem}\n"]
            md_lines.append(text_content)
            return "\n".join(md_lines)

    except FileNotFoundError:
        logger.info("libreoffice 未安装，跳过")
        return None
    except subprocess.TimeoutExpired:
        logger.error("libreoffice 转换超时")
        return None
    except Exception as e:
        logger.error(f"libreoffice 转换失败: {e}")
        return None


def _try_python_pptx(path: Path) -> Optional[str]:
    """使用 python-pptx 转换 PPTX 文件"""
    try:
        from pptx import Presentation

        prs = Presentation(path)
        md_lines = [f"# {path.stem}\n"]

        for slide_num, slide in enumerate(prs.slides, 1):
            md_lines.append(f"\n## 幻灯片 {slide_num}\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        md_lines.append(text)
                        md_lines.append("")

        return "\n".join(md_lines)

    except ImportError:
        logger.error("python-pptx 未安装，请运行: pip install python-pptx")
        return None
    except Exception as e:
        logger.error(f"PPTX 转换失败: {e}")
        return None


# ── 缓存管理 ─────────────────────────────────────────


def _converters_cache_dir() -> Path:
    from config import settings

    return Path(settings.cache_dir) / "converters"


def _content_hash(path: Path) -> str:
    """用文件内容 hash 做缓存 key"""
    data = path.read_bytes()
    return hashlib.sha256(data).hexdigest()[:16]


def _cache_path(path: Path) -> Path:
    return _converters_cache_dir() / f"{path.stem}_{_content_hash(path)}.md"


def _load_cache(path: Path) -> Optional[str]:
    """加载转换缓存，命中返回 Markdown 文本，未命中返回 None"""
    cache = _cache_path(path)
    if cache.exists():
        logger.info(f"[Converter] 缓存命中: {path.name}")
        return cache.read_text(encoding="utf-8")
    return None


def _save_cache(path: Path, content: str) -> None:
    """保存转换结果到缓存"""
    cache = _cache_path(path)
    cache.parent.mkdir(parents=True, exist_ok=True)
    cache.write_text(content, encoding="utf-8")
    logger.info(f"[Converter] 缓存已保存: {path.name}")
