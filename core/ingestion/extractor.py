"""
文件提取器

将文件转换为 Document 对象：

支持的格式及降级策略：
├── .md    → 直接读取
├── .pdf   → MinerU
├── .doc   → antiword → libreoffice 降级
├── .docx  → python-docx → XML 解析 → libreoffice 三级降级
├── .pptx  → python-pptx
└── .ppt   → 暂不支持

注意：分块由 chunker 负责，不在此处理。
"""

import hashlib
import subprocess
import sys
import tempfile
import uuid
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List, Optional

from loguru import logger

from core.model.models import Document

# ── MinerU 初始化 ──────────────────────────────────────────

MINERU_PATH = Path("/home/zjlab/Documents/build_LLMs/NLP_course_hf/MinerU")
MINERU_AVAILABLE = False

def _ensure_mineru():
    """延迟导入 MinerU（避免模块加载时就修改 sys.path）"""
    global MINERU_AVAILABLE, do_parse
    if MINERU_AVAILABLE:
        return
    if str(MINERU_PATH) not in sys.path:
        sys.path.insert(0, str(MINERU_PATH))
    try:
        from mineru.cli.client import do_parse as _do_parse
        do_parse = _do_parse
        MINERU_AVAILABLE = True
    except ImportError:
        MINERU_AVAILABLE = False


# ── 支持的格式 ─────────────────────────────────────────────


SUPPORTED_FORMATS = {".pdf", ".doc", ".docx", ".pptx", ".ppt", ".md"}


# ── 格式检测 ──────────────────────────────────────────────


def detect_format(file_path: str) -> str:
    """返回格式标识: pdf/doc/docx/pptx/ppt/md/unknown"""
    ext = Path(file_path).suffix.lower()
    format_map = {
        ".pdf": "pdf",
        ".doc": "doc",
        ".docx": "docx",
        ".pptx": "pptx",
        ".ppt": "ppt",
        ".md": "md",
    }
    return format_map.get(ext, "unknown")


# ── 缓存管理 ──────────────────────────────────────────────


def _cache_dir() -> Path:
    """获取转换缓存目录"""
    from config import settings

    return Path(settings.cache_dir) / "converters"


def _content_hash(path: Path) -> str:
    """用文件内容 hash 做缓存 key"""
    data = path.read_bytes()
    return hashlib.sha256(data).hexdigest()[:16]


def _cache_path(path: Path) -> Path:
    return _cache_dir() / f"{path.stem}_{_content_hash(path)}.md"


def _load_cache(path: Path) -> Optional[str]:
    """加载转换缓存，命中返回 Markdown 文本，未命中返回 None"""
    cache = _cache_path(path)
    if cache.exists():
        logger.info(f"[Extractor] 缓存命中: {path.name}")
        return cache.read_text(encoding="utf-8")
    return None


def _save_cache(path: Path, content: str) -> None:
    """保存转换结果到缓存"""
    cache = _cache_path(path)
    cache.parent.mkdir(parents=True, exist_ok=True)
    cache.write_text(content, encoding="utf-8")
    logger.info(f"[Extractor] 缓存已保存: {path.name}")


# ── 转换实现 ──────────────────────────────────────────────


def _convert_pdf(path: Path) -> str:
    """PDF → Markdown（使用 MinerU）"""
    _ensure_mineru()
    if not MINERU_AVAILABLE:
        raise RuntimeError("MinerU 不可用，请先安装 MinerU")

    # 检查缓存
    cached = _load_cache(path)
    if cached:
        return cached

    backend = "hybrid-auto-engine"

    with tempfile.TemporaryDirectory() as tmpdir:
        do_parse(
            output_dir=tmpdir,
            pdf_file_names=[path.stem],
            pdf_bytes_list=[path.read_bytes()],
            p_lang_list=["ch"],
            backend=backend,
            parse_method="auto",
            formula_enable=True,
            table_enable=True,
            f_draw_layout_bbox=False,
            f_draw_span_bbox=False,
            f_dump_md=True,
            f_dump_middle_json=True,
            f_dump_model_output=False,
            f_dump_orig_pdf=False,
            f_dump_content_list=False,
            f_make_md_mode="mm_markdown",
            start_page_id=0,
            end_page_id=None,
        )

        # 查找生成的 Markdown 文件
        backend_dir = backend.replace("-", "_").replace("_engine", "")
        md_file = Path(tmpdir) / path.stem / backend_dir / f"{path.stem}.md"

        if not md_file.exists():
            md_files = list(Path(tmpdir).rglob("*.md"))
            if md_files:
                md_file = md_files[0]
            else:
                raise RuntimeError("未找到生成的 Markdown 文件")

        md_content = md_file.read_text(encoding="utf-8")

    _save_cache(path, md_content)
    return md_content


def _convert_doc(path: Path) -> str:
    """DOC → Markdown: antiword → libreoffice 降级"""
    cached = _load_cache(path)
    if cached:
        return cached

    # 方法 1: antiword
    result = _try_antiword(path)
    if result:
        _save_cache(path, result)
        return result

    # 方法 2: libreoffice
    result = _convert_with_libreoffice(path)
    if result:
        _save_cache(path, result)
        return result

    raise RuntimeError(f"DOC 转换失败（antiword 和 libreoffice 均不可用）: {path.name}")


def _convert_docx(path: Path) -> str:
    """DOCX → Markdown: python-docx → XML 解析 → libreoffice 三级降级"""
    cached = _load_cache(path)
    if cached:
        return cached

    # 方法 1: python-docx
    result = _try_python_docx(path)
    if result:
        _save_cache(path, result)
        return result

    # 方法 2: 直接解析 XML
    result = _parse_docx_xml(path)
    if result:
        _save_cache(path, result)
        return result

    # 方法 3: libreoffice
    result = _convert_with_libreoffice(path)
    if result:
        _save_cache(path, result)
        return result

    raise RuntimeError(f"DOCX 转换失败（所有方法均不可用）: {path.name}")


def _convert_pptx(path: Path) -> str:
    """PPTX → Markdown: python-pptx"""
    cached = _load_cache(path)
    if cached:
        return cached

    result = _try_python_pptx(path)
    if result:
        _save_cache(path, result)
        return result

    raise RuntimeError(f"PPTX 转换失败: {path.name}")


# ── 具体转换实现 ───────────────────────────────────────────


def _try_antiword(path: Path) -> Optional[str]:
    """使用 antiword 转换 DOC 文件"""
    try:
        result = subprocess.run(
            ["antiword", "-m", "UTF-8", str(path)],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            md_lines = [f"# {path.stem}\n"]
            md_lines.append(result.stdout)
            return "\n".join(md_lines)
        logger.warning(f"antiword 转换失败: {result.stderr}")
    except FileNotFoundError:
        logger.info("antiword 未安装，跳过")
    except subprocess.TimeoutExpired:
        logger.warning("antiword 转换超时")
    except Exception as e:
        logger.warning(f"antiword 异常: {e}")
    return None


def _try_python_docx(path: Path) -> Optional[str]:
    """使用 python-docx 转换 DOCX 文件"""
    try:
        from docx import Document

        doc = Document(path)
        md_lines = [f"# {path.stem}\n"]

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                md_lines.append("")
                continue

            if para.style.name.startswith("Heading"):
                try:
                    level = int(para.style.name.replace("Heading ", "").strip())
                    md_lines.append(f"{'#' * (level + 1)} {text}")
                except ValueError:
                    md_lines.append(f"## {text}")
            else:
                md_lines.append(text)

        return "\n".join(md_lines)

    except ImportError:
        logger.info("python-docx 未安装，跳过")
        return None
    except Exception as e:
        logger.warning(f"python-docx 转换失败: {e}")
        return None


def _parse_docx_xml(path: Path) -> Optional[str]:
    """直接解析 word/document.xml（处理非标准 DOCX）"""
    try:
        with zipfile.ZipFile(path, "r") as zf:
            doc_xml = None
            for doc_path in ["word/document.xml", "word\\document.xml"]:
                try:
                    doc_xml = zf.read(doc_path)
                    break
                except KeyError:
                    continue

            if doc_xml is None:
                return None

            root = ET.fromstring(doc_xml)
            namespaces = {
                "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            }

            md_lines = [f"# {path.stem}\n"]
            for para in root.findall(".//w:p", namespaces):
                texts = [t.text or "" for t in para.findall(".//w:t", namespaces)]
                para_text = "".join(texts)
                if para_text.strip():
                    md_lines.append(para_text)
                else:
                    md_lines.append("")

            return "\n".join(md_lines)

    except Exception as e:
        logger.warning(f"XML 解析失败: {e}")
        return None


def _convert_with_libreoffice(path: Path) -> Optional[str]:
    """使用 libreoffice 将 DOC/DOCX 转换为纯文本"""
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)

            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "txt:Text (encoded)",
                    "--outdir",
                    str(temp_path),
                    str(path),
                ],
                capture_output=True,
                text=True,
                timeout=60,
            )

            if result.returncode != 0:
                logger.error(f"libreoffice 转换失败: {result.stderr}")
                return None

            txt_files = list(temp_path.glob("*.txt"))
            if not txt_files:
                logger.error("libreoffice 未生成输出文件")
                return None

            text_content = txt_files[0].read_text(encoding="utf-8", errors="ignore")
            md_lines = [f"# {path.stem}\n"]
            md_lines.append(text_content)
            return "\n".join(md_lines)

    except FileNotFoundError:
        logger.info("libreoffice 未安装，跳过")
        return None
    except subprocess.TimeoutExpired:
        logger.error("libreoffice 转换超时")
        return None
    except Exception as e:
        logger.error(f"libreoffice 转换失败: {e}")
        return None


def _try_python_pptx(path: Path) -> Optional[str]:
    """使用 python-pptx 转换 PPTX 文件"""
    try:
        from pptx import Presentation

        prs = Presentation(path)
        md_lines = [f"# {path.stem}\n"]

        for slide_num, slide in enumerate(prs.slides, 1):
            md_lines.append(f"\n## 幻灯片 {slide_num}\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        md_lines.append(text)
                        md_lines.append("")

        return "\n".join(md_lines)

    except ImportError:
        logger.error("python-pptx 未安装，请运行: pip install python-pptx")
        return None
    except Exception as e:
        logger.error(f"PPTX 转换失败: {e}")
        return None


# ── 主提取函数 ────────────────────────────────────────────


def convert_to_markdown(file_path: str) -> str:
    """
    统一入口：任意格式 → Markdown 文本

    Args:
        file_path: 文件路径

    Returns:
        Markdown 文本内容

    Raises:
        ValueError: 不支持的格式
        RuntimeError: 转换失败
    """
    path = Path(file_path)
    fmt = detect_format(path)

    if fmt == "unknown":
        raise ValueError(f"不支持的文件格式: {path.suffix}")

    if fmt == "md":
        return path.read_text(encoding="utf-8")

    if fmt == "pdf":
        return _convert_pdf(path)

    if fmt == "doc":
        return _convert_doc(path)

    if fmt == "docx":
        return _convert_docx(path)

    if fmt == "pptx":
        return _convert_pptx(path)

    if fmt == "ppt":
        raise RuntimeError("老 PPT 格式暂不支持，请转换为 PPTX")

    raise ValueError(f"未知的格式标识: {fmt}")


def _generate_doc_hash(content: str) -> str:
    """生成文档内容哈希"""
    return hashlib.sha256(content.encode()).hexdigest()[:16]


def extract(file_path: str) -> List[Document]:
    """
    提取文件内容为 Document 列表。

    输入：文件路径
    输出：List[Document]（整篇文档作为一个 Document，分块由 chunker 负责）

    Args:
        file_path: 文件路径

    Returns:
        Document 对象列表

    Raises:
        ValueError: 不支持的文件格式
        RuntimeError: 转换失败
    """
    path = Path(file_path)

    # 1. 格式检查
    fmt = detect_format(path)
    if fmt == "unknown":
        raise ValueError(f"不支持的文件格式: {path.suffix}")

    if fmt == "ppt":
        raise RuntimeError("老 PPT 格式暂不支持，请转换为 PPTX")

    # 2. 转换为 Markdown
    try:
        md_text = convert_to_markdown(str(path))
    except Exception as e:
        raise RuntimeError(f"文件转换失败: {file_path}, error: {e}")

    if not md_text.strip():
        logger.warning(f"[Extractor] 文件内容为空: {file_path}")
        return []

    # 3. 构建 Document
    doc_id = f"doc_{uuid.uuid4().hex[:8]}"
    doc_hash = _generate_doc_hash(md_text)

    doc = Document(
        content=md_text,
        metadata={
            "doc_id": doc_id,
            "doc_hash": doc_hash,
            "source": str(path),
            "source_format": fmt,
        },
    )

    logger.info(f"[Extractor] 提取完成: {file_path}, format={fmt}, doc_id={doc_id}")

    return [doc]


# ── MinerUPDFProcessor（保留兼容） ──────────────────────────────────────────


class MinerUPDFProcessor:
    """MinerU PDF 处理器 — 仅 PDF → Markdown（保留用于直接调用）"""

    def __init__(
        self,
        backend: str = "hybrid-auto-engine",
        model_source: str = "local",
        gpu_device: Optional[int] = None,
        parse_method: str = "auto",
    ):
        _ensure_mineru()
        if not MINERU_AVAILABLE:
            raise RuntimeError("MinerU 不可用，请先安装 MinerU")

        self.backend = backend
        self.model_source = model_source
        self.gpu_device = gpu_device
        self.parse_method = parse_method

    def process_pdf(self, pdf_path: str, output_dir: Optional[str] = None) -> str:
        """
        将 PDF 转换为 Markdown 并返回内容。

        Args:
            pdf_path: PDF 文件路径
            output_dir: 输出目录（默认为 PDF 同目录下的 output/）

        Returns:
            Markdown 文本内容
        """
        return _convert_pdf(Path(pdf_path))
